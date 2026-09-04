import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

shapes = []
for i in range(3):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1 + i*3.5), Inches(2), Inches(3), Inches(2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(20, 184, 166)
    tf = s.text_frame
    tf.text = f"Step {i+1}"
    shapes.append(s)

shape_ids = [s.shape_id for s in shapes]
child_nodes = []
c_tn_id = 3
for sp_id in shape_ids:
    node = f"""<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:cTn id="{c_tn_id}" fill="hold">
        <p:stCondLst>
          <p:cond delay="0"/>
        </p:stCondLst>
        <p:childTnLst>
          <p:par>
            <p:cTn id="{c_tn_id+1}" fill="hold">
              <p:stCondLst>
                <p:cond delay="0"/>
              </p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{c_tn_id+2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{c_tn_id+3}" dur="1" fill="hold"/>
                          <p:tgtEl>
                            <p:spTgt spid="{sp_id}"/>
                          </p:tgtEl>
                          <p:attrNameLst>
                            <p:attrName>style.visibility</p:attrName>
                          </p:attrNameLst>
                        </p:cBhvr>
                        <p:to>
                          <p:strVal val="visible"/>
                        </p:to>
                      </p:set>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>"""
    child_nodes.append(node)
    c_tn_id += 4

xml_str = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="always" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                {''.join(child_nodes)}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                  <p:spTgt spid="0"/>
                </p:tgtEl>
              </p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                  <p:spTgt spid="0"/>
                </p:tgtEl>
              </p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

timing_elm = parse_xml(xml_str)
slide._element.append(timing_elm)

prs.save("test_anim.pptx")
print("Saved test_anim.pptx with timing XML successfully")
