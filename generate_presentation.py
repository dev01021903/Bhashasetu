"""
MediDecode — AI-Powered Prescription & Health Companion
100% Guaranteed High-Contrast Visible Text Presentation Generator
Theme: Premium Modern Healthcare (Clean Light Canvas with Deep Blue & Teal Accents)
All text colors are explicitly stamped on runs (<a:rPr>) to eliminate any 'white on white' issues.
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation(output_path="d:\\translator\\MediDecode_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # -------------------------------------------------------------------------
    # High-Contrast Colors (Clean Light Canvas + Crisp Dark Navy Text)
    # -------------------------------------------------------------------------
    BG_CANVAS = RGBColor(248, 250, 252)     # Slate 50 (Crisp Light Background)
    BG_WHITE = RGBColor(255, 255, 255)      # Pure White for Cards
    BG_TEAL_TINT = RGBColor(240, 253, 250)  # Teal 50
    BG_BLUE_TINT = RGBColor(239, 246, 255)  # Blue 50
    BG_AMBER_TINT = RGBColor(254, 243, 199) # Amber 100
    BG_ROSE_TINT = RGBColor(255, 241, 242)  # Rose 50
    BG_GREEN_TINT = RGBColor(236, 253, 245) # Emerald 50

    # High Contrast Text Colors
    TEXT_DARK = RGBColor(15, 23, 42)        # Slate 900 (Nearly Black - Ultra Sharp)
    TEXT_NAVY = RGBColor(30, 58, 138)       # Blue 900
    TEXT_BODY = RGBColor(30, 41, 59)        # Slate 800 (Dark Charcoal for Body)
    TEXT_MUTED = RGBColor(71, 85, 105)      # Slate 600 (Dark Slate for Secondary text)
    TEXT_WHITE = RGBColor(255, 255, 255)    # White (ONLY used on solid dark backgrounds)

    # Accent Colors
    TEAL_PRIMARY = RGBColor(13, 148, 136)   # Teal 600
    BLUE_ACCENT = RGBColor(2, 132, 199)     # Sky 600
    EMERALD_GREEN = RGBColor(5, 150, 105)   # Emerald 600
    AMBER_ACCENT = RGBColor(217, 119, 6)    # Amber 600
    ROSE_ACCENT = RGBColor(225, 29, 72)     # Rose 600

    BORDER_SUBTLE = RGBColor(203, 213, 225) # Slate 300
    BORDER_TEAL = RGBColor(45, 212, 191)    # Teal 400
    BORDER_BLUE = RGBColor(147, 197, 253)   # Blue 300

    blank_layout = prs.slide_layouts[6]

    # Helper to add runs with 100% explicit styling
    def add_line(text_frame, text, font_size=12, color=TEXT_BODY, bold=False, italic=False, align=None):
        if not text_frame.paragraphs or not text_frame.paragraphs[0].text:
            p = text_frame.paragraphs[0]
            p.text = ""
        else:
            p = text_frame.add_paragraph()
            p.text = ""

        if align is not None:
            p.alignment = align

        run = p.add_run()
        run.text = text
        run.font.name = "Segoe UI"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return p, run

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CANVAS
        bg.line.fill.background()

        # Top Healthcare Accent Ribbon
        top_ribbon = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        )
        top_ribbon.fill.solid()
        top_ribbon.fill.fore_color.rgb = TEAL_PRIMARY
        top_ribbon.line.fill.background()

    def add_header(slide, tag_text, title_text, subtitle_text=""):
        # Tag pill
        tag_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(3.4), Inches(0.32)
        )
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = BG_TEAL_TINT
        tag_box.line.color.rgb = BORDER_TEAL
        tag_box.line.width = Pt(1)
        tf_tag = tag_box.text_frame
        tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_line(tf_tag, tag_text.upper(), font_size=9.5, color=TEAL_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

        # Main Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.70), Inches(11.733), Inches(0.85))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        add_line(tf_title, title_text, font_size=26, color=TEXT_DARK, bold=True)

        if subtitle_text:
            add_line(tf_title, subtitle_text, font_size=12.5, color=TEXT_MUTED, bold=False)

    def add_speaker_note(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Left Container
    left_box = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.1), Inches(7.5), Inches(5.6)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = BG_WHITE
    left_box.line.color.rgb = BORDER_SUBTLE
    left_box.line.width = Pt(1.5)

    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_l, "SIH 2026 / HEALTHCARE AI INNOVATION", font_size=11, color=TEAL_PRIMARY, bold=True)
    add_line(tf_l, "MediDecode", font_size=46, color=TEXT_DARK, bold=True)
    add_line(tf_l, "AI-Powered Prescription & Health Companion", font_size=18, color=BLUE_ACCENT, bold=True)
    add_line(tf_l, "“Decode. Understand. Remember. Stay Safe.”", font_size=15, color=TEAL_PRIMARY, italic=True)

    add_line(tf_l, "\n• AI Vision & Handwriting Decoder: Reads difficult clinic slips.", font_size=12, color=TEXT_BODY)
    add_line(tf_l, "• Vernacular Translation: Clear instructions in 11 Indian languages.", font_size=12, color=TEXT_BODY)
    add_line(tf_l, "• Food & Medication Safety: Proactive warnings on everyday diet clashes.", font_size=12, color=TEXT_BODY)
    add_line(tf_l, "• Smart Reminders & Adherence: Converts Latin frequencies into clock alarms.", font_size=12, color=TEXT_BODY)

    add_line(tf_l, "\nPresented by: Student Innovation Team | SIH 2026 Project Showcase", font_size=11.5, color=AMBER_ACCENT, bold=True)

    # Right Showcase Preview Card
    right_box = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.1), Inches(3.933), Inches(5.6)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = BG_TEAL_TINT
    right_box.line.color.rgb = BORDER_TEAL
    right_box.line.width = Pt(1.5)

    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_r, "📱 LIVE APP PREVIEW", font_size=12, color=TEAL_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_line(tf_r, "\n📷 Scanned Slip Decoded:", font_size=12, color=TEXT_DARK, bold=True)
    add_line(tf_r, "━━━━━━━━━━━━━━━━━━━", font_size=10, color=BORDER_TEAL)
    add_line(tf_r, "💊 Amoxicillin 500 mg", font_size=13, color=TEXT_DARK, bold=True)
    add_line(tf_r, "• Timing: After Meals", font_size=11, color=TEXT_BODY)
    add_line(tf_r, "• Frequency: Twice Daily (8 AM / 8 PM)", font_size=11, color=TEXT_BODY)
    add_line(tf_r, "• Confidence: 94% Verified 🟢", font_size=11, color=EMERALD_GREEN, bold=True)
    add_line(tf_r, "\n🥗 Food Safety Alert:", font_size=12, color=TEXT_DARK, bold=True)
    add_line(tf_r, "• Separate from dairy (milk/curd) by 2 hrs.", font_size=11, color=TEXT_BODY)
    add_line(tf_r, "\n🌐 Regional Language:", font_size=12, color=TEXT_DARK, bold=True)
    add_line(tf_r, "• भोजन के बाद दिन में दो बार लें।", font_size=11.5, color=TEXT_NAVY, bold=True)
    add_line(tf_r, "\n✔ Zero-Guessing Clinical Safety Rule", font_size=11, color=EMERALD_GREEN, bold=True, align=PP_ALIGN.CENTER)

    add_speaker_note(s1,
        "Good morning respected judges and mentors. Today, we are proud to present 'MediDecode' — an AI-Powered Prescription and Health Companion built for SIH 2026.\n\n"
        "Our mission is simple: 'Decode. Understand. Remember. Stay Safe.' Every day, millions of patients return from clinics unable to decipher doctor handwriting, confused by Latin abbreviations, or taking medicines with foods that neutralize their absorption. MediDecode solves this with a clinically-safe multimodal AI companion. Let's walk you through the problem, our architecture, and our product flow!"
    )

    # =========================================================================
    # SLIDE 2: THE PROBLEM
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Healthcare Challenge", "Healthcare Instructions Shouldn't Be Difficult to Understand", "Patients and caregivers struggle daily with complex, unreadable, and foreign medical documents.")

    problems = [
        ("✍️", "Difficult Handwritten Prescriptions", "Fast cursive shorthand, smudged inks, and illegible drug names create severe misinterpretation risks for patients and pharmacies.", ROSE_ACCENT, BG_ROSE_TINT),
        ("🧬", "Complex Medical Terminology", "Clinical jargon like 'Cholecystectomy', 'NPO', or 'Dyslipidemia' leaves patients completely in the dark about their actual diagnosis.", AMBER_ACCENT, BG_AMBER_TINT),
        ("🌐", "Severe Language Barriers", "Over 80% of prescriptions are written in English or Latin, creating huge communication gaps for non-English speaking families across India.", BLUE_ACCENT, BG_BLUE_TINT),
        ("⏰", "Forgotten Medicine Timings", "Confusing dosage frequencies (OD, BD, TID, HS) lead to skipped doses, double dosing, or prematurely stopping antibiotic courses.", TEAL_PRIMARY, BG_TEAL_TINT),
        ("🥗", "Confusion About Food & Medicine", "Patients unknowingly take antibiotics with milk (neutralizing absorption) or statins with grapefruit, triggering severe adverse effects.", AMBER_ACCENT, BG_AMBER_TINT),
        ("📑", "Dense Discharge Summaries", "Multi-page post-op hospital summaries overwhelm patients with clinical codes instead of simple, actionable home recovery instructions.", ROSE_ACCENT, BG_ROSE_TINT)
    ]

    p_w = Inches(3.72)
    p_h = Inches(2.3)
    p_xs = [Inches(0.8), Inches(4.8), Inches(8.8)]
    p_ys = [Inches(1.75), Inches(4.35)]

    for idx, (icon, title, desc, border_c, bg_c) in enumerate(problems):
        col = idx % 3
        row = idx // 3
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p_xs[col], p_ys[row], p_w, p_h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = border_c
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, f"{icon}  {title}", font_size=13.5, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{desc}", font_size=11, color=TEXT_BODY, bold=False)

    add_speaker_note(s2,
        "Here on Slide 2, we highlight the real crisis. Over 50% of chronic patients fail to follow instructions correctly. "
        "The problem is six-fold: 1) Doctors write under intense time pressure in dense cursive; 2) Medical terminology is unintelligible to laypeople; "
        "3) Language barriers exclude regional language speakers; 4) Medicine timings are forgotten; 5) Hidden food-drug clashes occur daily; "
        "and 6) Discharge summaries are too long and dense. This leads to avoidable hospital readmissions and antibiotic resistance."
    )

    # =========================================================================
    # SLIDE 3: OUR SOLUTION
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Our Innovation", "Meet MediDecode", "An end-to-end intelligent platform bridging clinical documents and daily recovery.")

    journey_steps = [
        ("01", "📷 Upload", "Prescription or Discharge Slip"),
        ("02", "🤖 AI Vision", "Multimodal OCR & Extraction"),
        ("03", "📋 Plan", "Structured Dosing & Clock Times"),
        ("04", "🍎 Food Safety", "Eat / Limit / Avoid Guidance"),
        ("05", "🌐 Vernacular", "11 Indian Native Languages"),
        ("06", "🔔 Reminders", "Smart Alarms & Dashboard")
    ]

    j_w = Inches(1.85)
    j_h = Inches(1.8)
    j_gap = Inches(0.12)
    j_start = Inches(0.8)

    for idx, (num, title, sub) in enumerate(journey_steps):
        x_pos = j_start + idx * (j_w + j_gap)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), j_w, j_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = TEAL_PRIMARY
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, f"STEP {num}", font_size=11, color=TEAL_PRIMARY, bold=True)
        add_line(tf, title, font_size=13, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{sub}", font_size=10, color=TEXT_MUTED, bold=False)

    val_cards = [
        ("🛡️ Strict Healthcare Safety", "Never invents or hallucinates medicine names. Unclear handwriting is flagged with confidence scores and user verification alerts.", TEAL_PRIMARY, BG_TEAL_TINT),
        ("🧠 Layperson Simplification", "'Explain this to me' feature translates surgical notes and clinical terms into 2-sentence conversational language.", BLUE_ACCENT, BG_BLUE_TINT),
        ("🤝 Informational Companion", "Clearly reminds patients to verify important medication decisions with licensed doctors and pharmacists.", EMERALD_GREEN, BG_GREEN_TINT)
    ]

    for idx, (title, desc, col, bg_c) in enumerate(val_cards):
        x_pos = Inches(0.8 + idx * 4.0)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.9), Inches(3.7), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = col
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, title, font_size=15, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{desc}", font_size=12, color=TEXT_BODY, bold=False)

    add_speaker_note(s3,
        "Slide 3 introduces MediDecode: Our platform is a complete patient journey from upload to daily reminder. "
        "Unlike generic consumer AI tools, MediDecode is built specifically for healthcare safety. "
        "Our core rule is ZERO GUESSING: If a doctor's cursive is smudged, our AI explicitly flags it in yellow/red with a verification alert. "
        "Notice our bottom three pillars: strict clinical guardrails, layperson simplification, and clear professional disclaimers."
    )

    # =========================================================================
    # SLIDE 4: HOW IT WORKS (5-STEP WORKFLOW)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "System Workflow", "How It Works: 5-Step Clinical Intelligence Flow", "From an unreadable paper slip to a personalized daily recovery companion.")

    steps5 = [
        ("01", "📷", "Upload Prescription", "User uploads image or PDF slip via mobile camera, file browser, or gallery.", BG_BLUE_TINT, BLUE_ACCENT),
        ("02", "👁️", "AI Vision Reads It", "Multimodal Gemini Vision scans handwriting, printed text, doctor stamps, and dates.", BG_TEAL_TINT, TEAL_PRIMARY),
        ("03", "💊", "Extract Medicines", "Parses generic name, strength, route, frequency, and exact course duration.", BG_WHITE, TEAL_PRIMARY),
        ("04", "🛡️", "Safety & Language Layer", "Cross-checks food interactions, translates directions, and evaluates confidence.", BG_AMBER_TINT, AMBER_ACCENT),
        ("05", "🔔", "Generate Reminders", "Constructs wall-clock dose alarms, food timing alerts, and adherence dashboard.", BG_GREEN_TINT, EMERALD_GREEN)
    ]

    card_w = Inches(2.15)
    card_h = Inches(4.3)
    start_x = Inches(0.8)
    gap = Inches(0.24)

    for idx, (num, icon, title, desc, bg_c, col) in enumerate(steps5):
        cx = start_x + idx * (card_w + gap)
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.85), card_w, card_h)
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, f"STEP {num}", font_size=11, color=col, bold=True)
        add_line(tf, f"\n{icon}", font_size=28, color=TEXT_DARK)
        add_line(tf, title, font_size=14, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{desc}", font_size=11, color=TEXT_BODY, bold=False)

        if idx < 4:
            ar_x = cx + card_w + Inches(0.04)
            ar = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ar_x, Inches(3.8), Inches(0.16), Inches(0.22))
            ar.fill.solid()
            ar.fill.fore_color.rgb = TEAL_PRIMARY
            ar.line.fill.background()

    ribbon = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.35), Inches(11.733), Inches(0.65))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BG_WHITE
    ribbon.line.color.rgb = BORDER_SUBTLE
    ribbon.line.width = Pt(1)
    tf_rib = ribbon.text_frame
    tf_rib.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_line(tf_rib, "PRODUCT LIFECYCLE:   UPLOAD  ➔  DECODE  ➔  UNDERSTAND  ➔  CHECK  ➔  REMIND  ➔  TRACK",
             font_size=12, color=TEAL_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    add_speaker_note(s4,
        "Slide 4 walks through our 5-step workflow: "
        "Step 1: Patient uploads the slip. "
        "Step 2: Gemini Vision multimodal OCR reads the handwritten prescription. "
        "Step 3: Clinical entities — medicine name, dosage, frequency — are extracted into structured JSON. "
        "Step 4: The safety layer evaluates confidence and checks drug-food interactions. "
        "Step 5: Daily reminder alarms are created and populated on the dashboard. "
        "At the bottom, you can see our product lifecycle in six words: Upload, Decode, Understand, Check, Remind, Track."
    )

    # =========================================================================
    # SLIDE 5: CORE FEATURES (8 FEATURE CARDS)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Platform Capabilities", "Core Features of MediDecode", "Engineered to deliver comprehensive clinical clarity and high user adoption.")

    features8 = [
        ("🔍", "AI Prescription Scanner", "Scans clinic slips, doctor notes, and stamped slips with specialized vision prompts.", TEAL_PRIMARY),
        ("💊", "Medicine Decoder", "Standardizes brand names into International Nonproprietary Names (INN) with strength & route.", BLUE_ACCENT),
        ("⏰", "Smart Medication Schedule", "Converts Latin shorthand (OD, BD, TID) into concrete clock alarms (8:00 AM, 8:00 PM).", EMERALD_GREEN),
        ("🔔", "Medicine Reminders", "Generates web push notifications and audio alerts so patients never skip a crucial dose.", AMBER_ACCENT),
        ("🍎", "Food & Medication Guidance", "Proactive dietary engine highlights foods to eat, limit, or avoid based on active medicines.", ROSE_ACCENT),
        ("⚠️", "Precaution & Safety Alerts", "Color-coded confidence scoring (Green/Yellow/Red) flags ambiguous cursive for pharmacist review.", AMBER_ACCENT),
        ("🌐", "Regional Language Translation", "Translates instructions into 11 Indian languages while preserving standardized medicine names.", BLUE_ACCENT),
        ("📄", "Discharge Summary Simplification", "'Explain this to me' feature translates complex surgical and hospital discharge notes.", TEAL_PRIMARY)
    ]

    f_w = Inches(2.78)
    f_h = Inches(2.3)
    f_xs = [Inches(0.8), Inches(3.78), Inches(6.76), Inches(9.74)]
    f_ys = [Inches(1.85), Inches(4.45)]

    for idx, (icon, title, desc, col) in enumerate(features8):
        c_col = idx % 4
        c_row = idx // 4
        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, f_xs[c_col], f_ys[c_row], f_w, f_h)
        c.fill.solid()
        c.fill.fore_color.rgb = BG_WHITE
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, f"{icon}  {title}", font_size=13, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{desc}", font_size=10.5, color=TEXT_BODY, bold=False)

    add_speaker_note(s5,
        "On Slide 5, we present our 8 core features. Notice how each feature targets a specific real-world friction point: "
        "from decoding cursive and converting Latin frequencies to 8:00 AM times, to our unique food guidance layer, "
        "11 Indian languages, and discharge summary simplification. Every module is designed for ease of use by non-technical patients."
    )

    # =========================================================================
    # SLIDE 6: FOOD & MEDICATION SAFETY (MAJOR DIFFERENTIATOR)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Clinical Differentiator", "Food & Medication Safety Engine", "A vital safety layer preventing common adverse diet-drug interactions.")

    top_flow = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.7))
    top_flow.fill.solid()
    top_flow.fill.fore_color.rgb = BG_TEAL_TINT
    top_flow.line.color.rgb = BORDER_TEAL
    top_flow.line.width = Pt(1)
    tf_tf = top_flow.text_frame
    tf_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(tf_tf, "FLOW:   [ Medicine Prescribed ]  ➔  [ Clinical Interaction Engine ]  ➔  [ Eat / Limit / Avoid Analysis ]  ➔  [ Safety Alert ]",
             font_size=12, color=TEAL_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    case_header = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.65), Inches(11.733), Inches(0.45))
    case_header.fill.solid()
    case_header.fill.fore_color.rgb = BG_WHITE
    case_header.line.color.rgb = BORDER_SUBTLE
    case_header.line.width = Pt(1)
    tf_ch = case_header.text_frame
    tf_ch.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(tf_ch, "REAL-WORLD CLINICAL EXAMPLE:  Ciprofloxacin 500 mg (Antibiotic)  |  Source: British National Formulary (BNF)",
             font_size=11, color=TEXT_NAVY, bold=True)

    cols_data = [
        ("🟢 EAT (SUITABLE)", "Foods that do not interfere:", [
            "• Plain cooked rice, steamed vegetables",
            "• Water, clear soups, herbal tea",
            "• Whole grain bread & lean proteins",
            "• High hydration promotes safe drug clearance"
        ], EMERALD_GREEN, BG_GREEN_TINT),
        ("🟡 LIMIT / TIME (CAUTION)", "Requires strict time separation:", [
            "• Milk, Curd, Yogurt, Cheese, Paneer",
            "• Calcium binds to antibiotic molecules",
            "• Reduces drug absorption by up to 70%",
            "• SAFE RULE: Separate dairy by 2 to 4 hours"
        ], AMBER_ACCENT, BG_AMBER_TINT),
        ("🔴 AVOID / CHECK (HIGH ALERT)", "Severe risks requiring attention:", [
            "• Calcium & Iron fortified supplements",
            "• Alcohol: Triggers severe stomach distress",
            "• Antacids containing aluminum/magnesium",
            "• Always verify with prescribing doctor"
        ], ROSE_ACCENT, BG_ROSE_TINT)
    ]

    for idx, (col_t, sub, bullets_list, col, bg_c) in enumerate(cols_data):
        x_pos = Inches(0.8 + idx * 4.0)
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.2), Inches(3.733), Inches(2.85))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = col
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, col_t, font_size=13, color=col, bold=True)
        add_line(tf, sub, font_size=10.5, color=TEXT_DARK, bold=True)

        for b in bullets_list:
            add_line(tf, b, font_size=10, color=TEXT_BODY, bold=False)

    disc_card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.18), Inches(11.733), Inches(0.55))
    disc_card.fill.solid()
    disc_card.fill.fore_color.rgb = BG_WHITE
    disc_card.line.color.rgb = BORDER_SUBTLE
    disc_card.line.width = Pt(1)
    tf_d = disc_card.text_frame
    tf_d.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(tf_d, "⚠️ CLINICAL SAFETY NOTICE: Always verify medication-specific food restrictions with a qualified doctor or pharmacist.",
             font_size=10.5, color=AMBER_ACCENT, bold=True, align=PP_ALIGN.CENTER)

    add_speaker_note(s6,
        "Judges, Slide 6 is one of MediDecode's biggest differentiators: our Food & Medication Safety Engine. "
        "Consider this common real-world scenario: A doctor prescribes Ciprofloxacin. The patient takes it with a glass of milk. "
        "The calcium in the milk binds to the drug, neutralizing absorption by 70%. The patient thinks the antibiotic failed, when in reality it was food chelation! "
        "MediDecode categorizes foods into Green (Suitable), Yellow (Caution / Time separation), and Red (Avoid). "
        "And we prominently display a safety disclaimer encouraging patients to verify decisions with their doctor."
    )

    # =========================================================================
    # SLIDE 7: SMART MEDICATION REMINDERS (REALISTIC PHONE UI)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Patient Experience", "Smart Medication Reminders & Daily UI", "A realistic, distraction-free smartphone experience that keeps patients on track.")

    # Smartphone Frame
    phone_frame = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.75), Inches(4.2), Inches(5.1))
    phone_frame.fill.solid()
    phone_frame.fill.fore_color.rgb = BG_WHITE
    phone_frame.line.color.rgb = TEXT_DARK
    phone_frame.line.width = Pt(2.5)

    tf_p = phone_frame.text_frame
    tf_p.word_wrap = True
    tf_p.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_p, "09:41 AM  ●●●  Wi-Fi  🔋 98%", font_size=9, color=TEXT_MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_line(tf_p, "\nGood Morning, Rahul 👋", font_size=14, color=TEXT_DARK, bold=True)
    add_line(tf_p, "Today's Schedule: 2 of 3 Doses Taken", font_size=11, color=TEXT_MUTED, bold=False)

    add_line(tf_p, "\n⏰ 09:00 AM — Amoxicillin 500 mg", font_size=11, color=TEXT_DARK, bold=True)
    add_line(tf_p, "• 1 tablet | After breakfast", font_size=10, color=TEXT_BODY)
    add_line(tf_p, "• STATUS: TAKEN ✓ (09:05 AM)", font_size=10, color=EMERALD_GREEN, bold=True)

    add_line(tf_p, "\n⏰ 02:00 PM — Pantoprazole 40 mg", font_size=11, color=TEXT_DARK, bold=True)
    add_line(tf_p, "• 1 tablet | Before lunch", font_size=10, color=TEXT_BODY)
    add_line(tf_p, "• STATUS: NEXT DOSE ⏱️ (In 45 mins)", font_size=10, color=BLUE_ACCENT, bold=True)

    add_line(tf_p, "\n⏰ 09:00 PM — Amoxicillin 500 mg", font_size=11, color=TEXT_DARK, bold=True)
    add_line(tf_p, "• 1 tablet | After dinner", font_size=10, color=TEXT_BODY)
    add_line(tf_p, "• STATUS: UPCOMING 📅", font_size=10, color=TEXT_MUTED, bold=True)

    states = [
        ("🟢 Taken ✓ (1-Tap Checkoff)", "Patients easily tap their doses. Confetti animations and adherence scores celebrate consistency.", EMERALD_GREEN, BG_GREEN_TINT),
        ("🔵 Next Dose (Smart Countdown)", "Highlights the immediate upcoming medicine with food instructions (e.g. 'Before lunch').", BLUE_ACCENT, BG_BLUE_TINT),
        ("⚪ Upcoming (Daily Overview)", "Shows evening and bedtime doses so patients and family caregivers can plan meals in advance.", TEAL_PRIMARY, BG_TEAL_TINT),
        ("🔴 Missed / Overdue Tracking", "If a dose is delayed, gentle reminder notifications prompt the patient and notify family members.", ROSE_ACCENT, BG_ROSE_TINT)
    ]

    for idx, (title, desc, col, bg_c) in enumerate(states):
        y_pos = Inches(1.8 + idx * 1.25)
        c = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), y_pos, Inches(6.733), Inches(1.1))
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, title, font_size=13, color=TEXT_DARK, bold=True)
        add_line(tf, desc, font_size=11, color=TEXT_BODY, bold=False)

    add_speaker_note(s7,
        "On Slide 7, we present the patient-facing Smartphone UI. "
        "Notice how simple and uncluttered the interface is: At 9:00 AM, the patient confirms their morning dose with a single tap. "
        "The system immediately indicates that their next dose is Pantoprazole at 2:00 PM, reminding them to take it before lunch. "
        "We support four clear states: Taken, Next Dose, Upcoming, and Missed. This distraction-free design empowers elderly patients to use the app independently."
    )

    # =========================================================================
    # SLIDE 8: REGIONAL LANGUAGE SUPPORT
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Vernacular Inclusion", "Regional Language Support", "Bridging the linguistic divide to make healthcare universally accessible.")

    quote_card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.95))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = BG_TEAL_TINT
    quote_card.line.color.rgb = TEAL_PRIMARY
    quote_card.line.width = Pt(1.5)
    tf_q = quote_card.text_frame
    tf_q.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(tf_q, "“Healthcare information becomes truly useful only when patients can understand it in their own language.”",
             font_size=14, color=TEAL_PRIMARY, bold=True, italic=True, align=PP_ALIGN.CENTER)

    lang_cards = [
        ("🇬🇧 English (Original)", "Standard Medical Direction", "Amoxicillin 500 mg\nTake 1 tablet twice daily after meals with water.", BLUE_ACCENT, BG_BLUE_TINT),
        ("🇮🇳 Hindi (हिन्दी)", "National Vernacular", "Amoxicillin 500 mg\nभोजन के बाद पानी के साथ दिन में दो बार 1 गोली लें।", TEAL_PRIMARY, BG_TEAL_TINT),
        ("🇮🇳 Kannada (ಕನ್ನಡ)", "Regional State Language", "Amoxicillin 500 mg\nಊಟದ ನಂತರ ದಿನಕ್ಕೆ ಎರಡು ಬಾರಿ ನೀರಿನೊಂದಿಗೆ 1 ಮಾತ್ರೆ ತೆಗೆದುಕೊಳ್ಳಿ.", EMERALD_GREEN, BG_GREEN_TINT)
    ]

    for idx, (lang_t, sub, text_val, col, bg_c) in enumerate(lang_cards):
        x_pos = Inches(0.8 + idx * 4.0)
        c = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.95), Inches(3.733), Inches(2.55))
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, lang_t, font_size=14, color=TEXT_DARK, bold=True)
        add_line(tf, sub, font_size=11, color=col, bold=True)
        add_line(tf, f"\n{text_val}", font_size=12, color=TEXT_DARK, bold=True)

    diff_card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.7), Inches(11.733), Inches(1.1))
    diff_card.fill.solid()
    diff_card.fill.fore_color.rgb = BG_WHITE
    diff_card.line.color.rgb = BORDER_SUBTLE
    diff_card.line.width = Pt(1)
    tf_diff = diff_card.text_frame
    tf_diff.vertical_anchor = MSO_ANCHOR.TOP
    add_line(tf_diff, "🔑 KEY CLINICAL SAFETY RULE IN TRANSLATION:", font_size=11.5, color=AMBER_ACCENT, bold=True)
    add_line(tf_diff, "Medicine names (e.g. 'Amoxicillin 500 mg') are preserved in Roman script to prevent dangerous drug confusion at local pharmacy counters, while frequencies and food timings translate seamlessly into the patient's native script.",
             font_size=11, color=TEXT_BODY, bold=False)

    add_speaker_note(s8,
        "Judges, Slide 8 addresses India's linguistic diversity. "
        "A prescription in English is ineffective if an elderly patient in Karnataka or Uttar Pradesh cannot understand whether to take it before or after food. "
        "Here, we demonstrate the exact same instruction translated into English, Hindi, and Kannada. "
        "Crucially, notice our safety standard: Standardized INN medicine names remain in Roman script, while the timing, dose, and food directions translate into the native script. "
        "This ensures the local chemist always knows the exact chemical entity while the patient understands the directions effortlessly."
    )

    # =========================================================================
    # SLIDE 9: TECHNOLOGY ARCHITECTURE
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Technical Blueprint", "Technology Architecture & Data Pipeline", "Decoupled full-stack architecture built for high reliability and patient privacy.")

    pipe_steps = [
        ("1. User", "Patient / Caregiver"),
        ("2. Upload", "Camera / File Picker"),
        ("3. Vision OCR", "Gemini Vision Multimodal"),
        ("4. Extraction", "Medicine, Dose, Timing"),
        ("5. Safety Check", "Curated BNF / FDA DB"),
        ("6. Translation", "11 Native Languages"),
        ("7. Reminders", "Notification Engine"),
        ("8. Dashboard", "Daily Adherence Portal")
    ]

    p_w = Inches(1.36)
    p_h = Inches(1.6)
    p_gap = Inches(0.12)
    p_start = Inches(0.8)

    for idx, (t, s) in enumerate(pipe_steps):
        x_pos = p_start + idx * (p_w + p_gap)
        c = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), p_w, p_h)
        c.fill.solid()
        c.fill.fore_color.rgb = BG_WHITE
        c.line.color.rgb = TEAL_PRIMARY
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, t, font_size=10.5, color=TEAL_PRIMARY, bold=True)
        add_line(tf, f"\n{s}", font_size=9.5, color=TEXT_DARK, bold=True)

    arch_cols = [
        ("Frontend & Presentation Layer", "• Next.js 15 (App Router) & React 19\n• TypeScript & Tailwind CSS\n• Responsive Mobile-First PWA\n• Lucide Medical Iconography\n• Status: IMPLEMENTED & TESTED ✔️", TEAL_PRIMARY, BG_TEAL_TINT),
        ("AI & Clinical Knowledge Engine", "• Google Gemini Vision Multimodal API\n• Curated BNF / FDA Food Interaction DB\n• 11-Language Medical Dictionary\n• Zero-Guessing Verification Prompts\n• Status: IMPLEMENTED & TESTED ✔️", BLUE_ACCENT, BG_BLUE_TINT),
        ("Backend, Security & Roadmap", "• Supabase PostgreSQL with Row Level Security\n• Next.js Edge API Route Handlers\n• Web Notification Reminders Engine\n• Planned: Flutter Native & ABDM EMR Sync\n• Status: CORE LIVE / EXPANDING 🚀", EMERALD_GREEN, BG_GREEN_TINT)
    ]

    for idx, (title, points, col, bg_c) in enumerate(arch_cols):
        x_pos = Inches(0.8 + idx * 4.0)
        c = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.6), Inches(3.733), Inches(3.2))
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        add_line(tf, title, font_size=13, color=TEXT_DARK, bold=True)
        add_line(tf, f"\n{points}", font_size=10.5, color=TEXT_BODY, bold=False)

    add_speaker_note(s9,
        "On Slide 9, we review our technical architecture. "
        "At the top is our end-to-end data pipeline: from user upload to Gemini Vision multimodal extraction, "
        "safety verification against our curated BNF database, multi-lingual simplification, and browser reminder dispatch. "
        "Notice our technology stack: We built with Next.js 15, React 19, TypeScript, and Supabase PostgreSQL with Row Level Security for patient privacy. "
        "We are transparent: Our core AI, scheduling, and food engines are live and tested, with Flutter and ABDM integration on our roadmap."
    )

    # =========================================================================
    # SLIDE 10: WHY MEDIDECODE? (COMPARISON MATRIX)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Competitive Advantage", "Why MediDecode? Traditional Approach vs MediDecode", "Transforming passive, confusing paper slips into an active, verified digital recovery companion.")

    left_comp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.1))
    left_comp.fill.solid()
    left_comp.fill.fore_color.rgb = BG_ROSE_TINT
    left_comp.line.color.rgb = ROSE_ACCENT
    left_comp.line.width = Pt(1.5)

    tf_lc = left_comp.text_frame
    tf_lc.word_wrap = True
    tf_lc.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_lc, "❌  TRADITIONAL HEALTHCARE APPROACH", font_size=14, color=ROSE_ACCENT, bold=True)

    trad_points = [
        ("\n• Difficult Cursive Handwriting:", "Misread drug names and wrong dosages leading to medical emergencies."),
        ("• Complex Latin Abbreviations:", "Terms like 'BD', 'TID', 'HS' confuse patients about exact clock timings."),
        ("• English-Only Language Barrier:", "Non-English speakers cannot understand their own prescribed medications."),
        ("• Manual Medicine Tracking:", "Patients rely on memory, leading to missed doses or early course termination."),
        ("• Zero Dietary Awareness:", "Patients take medicines with calcium or citrus without knowing it blocks absorption."),
        ("• Dense Discharge Papers:", "Multi-page technical hospital documents that patients cannot interpret at home.")
    ]
    for pt, desc in trad_points:
        add_line(tf_lc, f"{pt} {desc}", font_size=10.5, color=TEXT_DARK, bold=False)

    right_comp = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))
    right_comp.fill.solid()
    right_comp.fill.fore_color.rgb = BG_GREEN_TINT
    right_comp.line.color.rgb = EMERALD_GREEN
    right_comp.line.width = Pt(1.5)

    tf_rc = right_comp.text_frame
    tf_rc.word_wrap = True
    tf_rc.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_rc, "✔️  THE MEDIDECODE PLATFORM", font_size=14, color=EMERALD_GREEN, bold=True)

    medi_points = [
        ("\n• AI-Assisted Multimodal OCR:", "Reads printed and cursive text with confidence badges and zero-guessing rules."),
        ("• Plain-Language Clock Timelines:", "Converts 'BD after food' into concrete 8:00 AM and 8:00 PM alarm slots."),
        ("• 11 Regional Indian Languages:", "Instructions in Hindi, Bengali, Kannada, Tamil, etc., while preserving INN drug names."),
        ("• Smart Automated Reminders:", "Web push notifications and 1-tap taken checkoff drive daily treatment adherence."),
        ("• Curated Food & Safety Engine:", "Proactive alerts on everyday foods to Eat, Limit, or Avoid based on verified BNF data."),
        ("• 'Explain This to Me' Summaries:", "Translates complicated post-op hospital records into 2-sentence actionable advice.")
    ]
    for pt, desc in medi_points:
        add_line(tf_rc, f"{pt} {desc}", font_size=10.5, color=TEXT_DARK, bold=False)

    add_speaker_note(s10,
        "On Slide 10, we contrast the Traditional Approach with MediDecode. "
        "The traditional method is passive and error-prone: difficult handwriting, obscure Latin abbreviations, "
        "English barriers, forgotten doses, food clashes, and dense discharge documents. "
        "MediDecode replaces this entire cycle with active digital intelligence: multimodal OCR with confidence scoring, "
        "plain-language clock times, 11 Indian languages, smart push reminders, curated food safety guidance, and 2-sentence discharge summaries. "
        "It transforms patient adherence from anxiety to clarity."
    )

    # =========================================================================
    # SLIDE 11: IMPACT
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "Social Value & Reach", "Real-World Impact & Beneficiaries", "Empowering patients, families, and healthcare systems across India.")

    impact_banner = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.95))
    impact_banner.fill.solid()
    impact_banner.fill.fore_color.rgb = BG_WHITE
    impact_banner.line.color.rgb = TEAL_PRIMARY
    impact_banner.line.width = Pt(1.5)
    tf_ib = impact_banner.text_frame
    tf_ib.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_line(tf_ib, "“MediDecode helps patients understand and follow healthcare instructions more confidently.”",
             font_size=14, color=TEAL_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    beneficiaries = [
        ("👵 Elderly Patients", "Large fonts, simple alarm clocks, and zero confusing medical jargon make daily medicine management effortless.", TEAL_PRIMARY),
        ("👨‍👩‍👧 Families & Caregivers", "Enables children and relatives to monitor their elderly parents' doses and recovery schedules remotely.", BLUE_ACCENT),
        ("💊 Polypharmacy Patients", "Chronic patients managing 5+ daily medications get an organized, non-conflicting master timetable.", EMERALD_GREEN),
        ("🇮🇳 Regional-Language Users", "Removes English literacy barriers, democratizing healthcare accessibility across Tier 2, 3, and rural India.", AMBER_ACCENT),
        ("🏥 Hospital Discharge Patients", "Clear post-surgery wound care and dietary guidelines prevent avoidable emergency readmissions.", ROSE_ACCENT)
    ]

    for idx, (title, desc, col) in enumerate(beneficiaries):
        y_pos = Inches(2.95 + idx * 0.82)
        c = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(0.72))
        c.fill.solid()
        c.fill.fore_color.rgb = BG_WHITE
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p, r = add_line(tf, f"{title}:  ", font_size=12, color=col, bold=True)
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT_DARK

    add_speaker_note(s11,
        "Slide 11 highlights our social value and community impact. "
        "Our mission is simple: 'MediDecode helps patients understand and follow healthcare instructions more confidently.' "
        "We empower five key groups: Elderly citizens who struggle with small print; family caregivers managing elderly parents from a distance; "
        "polypharmacy patients with multiple prescriptions; vernacular users across India; and post-surgical discharge patients who need clear home guidelines. "
        "This is healthcare technology with tangible human value."
    )

    # =========================================================================
    # SLIDE 12: FUTURE SCOPE & CLOSING
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_header(s12, "Roadmap & Vision", "Future Scope & Vision for MediDecode", "Scaling MediDecode into a comprehensive universal health companion.")

    left_road = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.8), Inches(5.1))
    left_road.fill.solid()
    left_road.fill.fore_color.rgb = BG_WHITE
    left_road.line.color.rgb = BORDER_SUBTLE
    left_road.line.width = Pt(1.5)

    tf_lr = left_road.text_frame
    tf_lr.word_wrap = True
    tf_lr.vertical_anchor = MSO_ANCHOR.TOP

    add_line(tf_lr, "🚀  FUTURE INNOVATION ROADMAP", font_size=14, color=TEAL_PRIMARY, bold=True)

    roadmap_items = [
        ("\n• Voice Assistant & Audio Readout:", "Natural voice instructions in 11 languages for illiterate or visually impaired patients."),
        ("• 22 Official Indian Languages:", "Expanding vernacular dialect coverage across all Indian states."),
        ("• Doctor & Pharmacist Verification Portal:", "Enables local chemists to verify flagged cursive medicines in 1-click."),
        ("• Dedicated Caregiver Companion App:", "Real-time alerts to family members if an elderly patient misses an insulin or BP dose."),
        ("• Wearable & Smartwatch Integration:", "Haptic wrist buzzers for discrete, timely medication alerts on smart bands."),
        ("• ABDM & EMR Integration:", "Seamless sync with Ayushman Bharat Digital Mission and electronic health records.")
    ]
    for r_title, r_desc in roadmap_items:
        add_line(tf_lr, f"{r_title} {r_desc}", font_size=10.5, color=TEXT_DARK, bold=False)

    right_close = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.8), Inches(4.633), Inches(5.1))
    right_close.fill.solid()
    right_close.fill.fore_color.rgb = BG_TEAL_TINT
    right_close.line.color.rgb = TEAL_PRIMARY
    right_close.line.width = Pt(2)

    tf_rc = right_close.text_frame
    tf_rc.word_wrap = True
    tf_rc.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_line(tf_rc, "🩺\nMediDecode", font_size=36, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_line(tf_rc, "“Decode. Understand. Remember. Stay Safe.”", font_size=13.5, color=TEAL_PRIMARY, bold=True, italic=True, align=PP_ALIGN.CENTER)
    add_line(tf_rc, "\nThank you respected judges and professors for your time and guidance.", font_size=11.5, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_line(tf_rc, "\nWe are now open for\nQuestions & Answers / Live Demo!", font_size=15, color=AMBER_ACCENT, bold=True, align=PP_ALIGN.CENTER)

    add_speaker_note(s12,
        "In conclusion, on Slide 12: "
        "Our roadmap for MediDecode includes native voice readout for visually impaired patients, expanding to 22 Indian languages, "
        "smartwatch haptic alarms, and seamless integration with the Ayushman Bharat Digital Mission (ABDM). "
        "MediDecode is not just an academic project — it is a production-ready, safety-first healthcare companion. "
        "Remember our tagline: 'Decode. Understand. Remember. Stay Safe.' "
        "Thank you so much, and we are now excited to demonstrate our live product and answer your questions!"
    )

    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

    # Update both Desktop and Downloads
    destinations = [
        os.path.expanduser("~/OneDrive/Desktop/MediDecode_Presentation.pptx"),
        os.path.expanduser("~/Desktop/MediDecode_Presentation.pptx"),
        os.path.expanduser("~/Downloads/MediDecode_Presentation.pptx")
    ]
    for d in destinations:
        try:
            os.makedirs(os.path.dirname(d), exist_ok=True)
            shutil.copy2(output_path, d)
            print(f"Updated copy at: {d}")
        except Exception as e:
            pass

if __name__ == "__main__":
    create_presentation("d:\\translator\\MediDecode_Presentation.pptx")
